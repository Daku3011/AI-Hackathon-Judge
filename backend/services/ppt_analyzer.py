from pptx import Presentation
import io

def extract_text_from_ppt(file_content: bytes) -> str:
    """
    Extracts text from a PPT/PPTX file content.
    Returns a string containing all text found in the slides.
    """
    try:
        ppt_file = io.BytesIO(file_content)
        prs = Presentation(ppt_file)
        
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                             text_runs.append(run.text)
                             
        return "\n".join(text_runs)
    except Exception as e:
        print(f"Error extracting text from PPT: {e}")
        return ""
